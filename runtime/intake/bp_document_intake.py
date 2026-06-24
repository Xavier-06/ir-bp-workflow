#!/usr/bin/env python3
"""
BP Document Intake — 真实逻辑

1. 接收文件路径（PDF/PPTX/图片/Word）
2. OCR/VL 模型解析文档内容
3. 结构化抽取：公司名、行业、融资阶段、商业模式、团队、财务数据
4. 输出 bp_ocr_text.txt + bp_step0_profile.json

使用小马算力 qwen3-vl-30b-a3b-instruct 做 OCR 和结构化抽取。
"""
from __future__ import annotations

import base64
import importlib.util
import json
import mimetypes
import os
import re
import shutil
import subprocess
import time
import zipfile
from pathlib import Path
from typing import Any, Optional

import requests

# ── 配置 ──────────────────────────────────────────────
ROOT = Path(__file__).resolve().parents[2]
DEFAULT_VL_API_BASE = ""  # Set via VL_API_BASE env var
DEFAULT_VL_MODEL = "qwen3-vl-30b-a3b-instruct"


def _read_env_file() -> dict[str, str]:
    """读取项目 .credentials 中的 OCR/VL 配置。"""
    cred_path = ROOT / ".credentials" / "investment-research.env"
    values: dict[str, str] = {}
    if not cred_path.exists():
        return values
    for raw_line in cred_path.read_text(encoding="utf-8").splitlines():
        line = raw_line.strip()
        if not line or line.startswith("#") or "=" not in line:
            continue
        key, value = line.split("=", 1)
        values[key.strip()] = value.strip().strip("'\"")
    return values


def _resolve_vl_config() -> tuple[str, str, str]:
    """解析 VL 调用配置，兼容历史 VL_* 和 BP_OCR_* 两套变量。"""
    file_values = _read_env_file()
    api_base = (
        os.environ.get("VL_API_BASE")
        or os.environ.get("BP_OCR_BASE_URL")
        or file_values.get("VL_API_BASE")
        or file_values.get("BP_OCR_BASE_URL")
        or DEFAULT_VL_API_BASE
    )
    api_key = (
        os.environ.get("VL_API_KEY")
        or os.environ.get("BP_OCR_API_KEY")
        or os.environ.get("DASHSCOPE_API_KEY")
        or file_values.get("VL_API_KEY")
        or file_values.get("BP_OCR_API_KEY")
        or file_values.get("DASHSCOPE_API_KEY")
        or ""
    )
    model = (
        os.environ.get("VL_MODEL")
        or os.environ.get("BP_OCR_MODEL")
        or file_values.get("VL_MODEL")
        or file_values.get("BP_OCR_MODEL")
        or DEFAULT_VL_MODEL
    )
    return api_base.rstrip("/"), api_key, model

# 支持的文件类型
SUPPORTED_EXTENSIONS = {".pdf", ".pptx", ".ppt", ".docx", ".doc",
                        ".png", ".jpg", ".jpeg", ".bmp", ".tiff", ".webp"}


# ── 文件处理 ──────────────────────────────────────────

def _file_to_base64(file_path: Path) -> tuple[str, str]:
    """读取文件并转 base64，返回 (base64_str, mime_type)"""
    mime, _ = mimetypes.guess_type(str(file_path))
    if mime is None:
        # 默认用 application/octet-stream
        mime = "application/octet-stream"
    with open(file_path, "rb") as f:
        b64 = base64.b64encode(f.read()).decode("utf-8")
    return b64, mime


def _extract_images_from_pptx(pptx_path: Path, output_dir: Path) -> list[Path]:
    """从 PPTX 中提取图片（PPTX 是 ZIP 格式）"""
    images = []
    output_dir.mkdir(parents=True, exist_ok=True)
    with zipfile.ZipFile(pptx_path, 'r') as z:
        for name in z.namelist():
            if name.startswith('ppt/media/') and any(
                name.lower().endswith(ext) for ext in ('.png', '.jpg', '.jpeg', '.bmp', '.tiff', '.webp')
            ):
                data = z.read(name)
                img_name = Path(name).name
                img_path = output_dir / img_name
                img_path.write_bytes(data)
                images.append(img_path)
    return images


def _render_pptx_with_libreoffice(pptx_path: Path, work_dir: Path) -> Optional[Path]:
    """用 LibreOffice 把 PPTX/PPT 渲染为逐页 PNG。
    
    流程：PPTX → PDF（LibreOffice headless）→ 逐页 PNG（pdftoppm）
    返回 PNG 目录路径，失败返回 None。
    """
    import subprocess
    
    # 查找 LibreOffice
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice:
        # macOS 常见路径
        mac_path = Path("/Applications/LibreOffice.app/Contents/MacOS/soffice")
        if mac_path.exists():
            soffice = str(mac_path)
    if not soffice:
        return None
    
    # 输出目录
    render_dir = work_dir / "rendered_slides"
    render_dir.mkdir(parents=True, exist_ok=True)
    
    # Step 1: PPTX → PDF
    try:
        result = subprocess.run(
            [soffice, "--headless", "--convert-to", "pdf",
             "--outdir", str(render_dir), str(pptx_path)],
            capture_output=True, text=True, timeout=120
        )
        if result.returncode != 0:
            return None
    except Exception:
        return None
    
    # 找到生成的 PDF
    pdf_files = list(render_dir.glob("*.pdf"))
    if not pdf_files:
        return None
    pdf_path = pdf_files[0]
    
    # Step 2: PDF → 逐页 PNG（pdftoppm）
    png_dir = work_dir / "slide_pngs"
    png_dir.mkdir(parents=True, exist_ok=True)
    
    pdftoppm = shutil.which("pdftoppm")
    if pdftoppm:
        try:
            prefix = str(png_dir / "slide")
            result = subprocess.run(
                [pdftoppm, "-png", "-r", "200", str(pdf_path), prefix],
                capture_output=True, text=True, timeout=120
            )
            if result.returncode == 0 and list(png_dir.glob("*.png")):
                # 清理 PDF
                pdf_path.unlink(missing_ok=True)
                return png_dir
        except Exception:
            pass
    
    # Fallback: pdf2image Python 包
    try:
        from pdf2image import convert_from_path
        images = convert_from_path(str(pdf_path), dpi=200)
        for i, img in enumerate(images, 1):
            out_path = png_dir / f"slide-{i:02d}.png"
            img.save(str(out_path), "PNG")
        pdf_path.unlink(missing_ok=True)
        return png_dir
    except ImportError:
        pass
    except Exception:
        pass
    
    # 都失败了，清理
    pdf_path.unlink(missing_ok=True)
    return None


def _extract_images_from_docx(docx_path: Path, output_dir: Path) -> list[Path]:
    """从 DOCX 中提取图片"""
    images = []
    output_dir.mkdir(parents=True, exist_ok=True)
    with zipfile.ZipFile(docx_path, 'r') as z:
        for name in z.namelist():
            if name.startswith('word/media/') and any(
                name.lower().endswith(ext) for ext in ('.png', '.jpg', '.jpeg', '.bmp', '.tiff', '.webp')
            ):
                data = z.read(name)
                img_name = Path(name).name
                img_path = output_dir / img_name
                img_path.write_bytes(data)
                images.append(img_path)
    return images


# ── VL 模型调用 ───────────────────────────────────────

def _vl_chat(messages: list[dict], max_tokens: int = 4096) -> str:
    """调用 VL 模型的 chat completions 接口"""
    vl_api_base, vl_api_key, vl_model = _resolve_vl_config()
    if not vl_api_key:
        raise RuntimeError("VL_API_KEY/BP_OCR_API_KEY 未配置，无法调用 BP OCR 模型")
    url = f"{vl_api_base}/chat/completions"
    headers = {
        "Authorization": f"Bearer {vl_api_key}",
        "Content-Type": "application/json",
    }
    payload = {
        "model": vl_model,
        "messages": messages,
        "max_tokens": max_tokens,
    }
    resp = requests.post(url, headers=headers, json=payload, timeout=120)
    resp.raise_for_status()
    data = resp.json()
    return data["choices"][0]["message"]["content"]


def _ocr_image(image_path: Path, page_hint: str = "") -> str:
    """用 VL 模型 OCR 单张图片"""
    b64, mime = _file_to_base64(image_path)
    content = []
    if page_hint:
        content.append({"type": "text", "text": f"这是文档的{page_hint}。"})
    content.append({
        "type": "image_url",
        "image_url": {"url": f"data:{mime};base64,{b64}"},
    })
    content.append({
        "type": "text",
        "text": (
            "请完整提取这张图片中的所有文字内容，保持原始排版结构。"
            "如果这是团队介绍页、顾问页、组织架构页或人物介绍页，必须逐行保留每个人的姓名、职务、头衔、教育背景、工作经历、所属机构、奖项、时间、百分比和括号内说明。"
            "不要把人物信息概括成一句摘要，不要漏掉图片中的小字、角标、表格、页脚、图注。"
            "只输出识别到的原始文字，不要解释，不要总结。"
        ),
    })
    messages = [{"role": "user", "content": content}]
    return _vl_chat(messages, max_tokens=4096)


def _ocr_pdf_image_paths(image_paths: list[Path], backend: str) -> tuple[str, int]:
    page_texts: list[str] = []
    total_pages = len(image_paths)
    print(f"  🔍 {backend} + 小马算力 VL OCR: {total_pages} 页", flush=True)
    for i, image_path in enumerate(image_paths, 1):
        try:
            t = _ocr_image(image_path, page_hint=f"第{i}页")
            if t and len(t) > 10:
                page_texts.append(f"--- Page {i} ---\n{t}")
            else:
                page_texts.append(f"--- Page {i} ---\n（VL识别无文字）")
            if i % 5 == 0 or i == total_pages:
                print(f"    OCR 进度: {i}/{total_pages}", flush=True)
        except Exception as e:
            page_texts.append(f"--- Page {i} ---\n[OCR失败: {e}]")
        time.sleep(0.3)
    return "\n\n".join(page_texts), total_pages


def _render_pdf_with_pdftoppm(pdf_path: Path, output_dir: Path, dpi: int = 200) -> list[Path]:
    pdftoppm = shutil.which("pdftoppm")
    if not pdftoppm:
        return []
    page_dir = output_dir / "pdf_pages"
    page_dir.mkdir(parents=True, exist_ok=True)
    prefix = str(page_dir / "page")
    result = subprocess.run(
        [pdftoppm, "-png", "-r", str(dpi), str(pdf_path), prefix],
        capture_output=True,
        text=True,
        timeout=180,
    )
    if result.returncode != 0:
        print(f"  ⚠ pdftoppm 渲染失败: {result.stderr or result.stdout}", flush=True)
        return []
    return sorted(page_dir.glob("page-*.png"))


def _render_pdf_with_pdf2image(pdf_path: Path, output_dir: Path, dpi: int = 200) -> list[Path]:
    from pdf2image import convert_from_path
    page_dir = output_dir / "pdf2image_pages"
    page_dir.mkdir(parents=True, exist_ok=True)
    images = convert_from_path(str(pdf_path), dpi=dpi, thread_count=4)
    paths: list[Path] = []
    for i, img in enumerate(images, 1):
        path = page_dir / f"page-{i}.jpg"
        img.save(path, format="JPEG", quality=90)
        paths.append(path)
    return paths


def _ocr_pdf(pdf_path: Path, output_dir: Path) -> tuple[str, int]:
    """OCR PDF — 主链必须逐页渲染为图片后调用小马算力 VL。"""

    # 方法 1（首选）: 系统 pdftoppm 逐页转 PNG → 小马算力 VL OCR。
    try:
        image_paths = _render_pdf_with_pdftoppm(pdf_path, output_dir, dpi=200)
        if image_paths:
            return _ocr_pdf_image_paths(image_paths, backend="pdftoppm")
    except Exception as e:
        print(f"  ⚠ pdftoppm + VL OCR 失败 ({e})，尝试 pdf2image", flush=True)

    # 方法 2: pdf2image 逐页转图片 → 小马算力 VL OCR。
    try:
        print("  📄 PDF 逐页渲染 (pdf2image) ...", flush=True)
        image_paths = _render_pdf_with_pdf2image(pdf_path, output_dir, dpi=200)
        if image_paths:
            return _ocr_pdf_image_paths(image_paths, backend="pdf2image")
    except ImportError:
        print("  ⚠ pdf2image 未安装，跳过 pdf2image VL OCR", flush=True)
    except Exception as e:
        print(f"  ⚠ pdf2image + VL OCR 失败 ({e})，回退", flush=True)

    # 方法 3: 整份 PDF base64 直接传 VL（可能只读首页）
    try:
        b64, mime = _file_to_base64(pdf_path)
        messages = [{
            "role": "user",
            "content": [
                {"type": "image_url", "image_url": {"url": f"data:application/pdf;base64,{b64}"}},
                {"type": "text", "text": "请完整提取这份文档中的所有文字内容，保持原始排版结构。只输出文字，不要解释。"},
            ],
        }]
        result = _vl_chat(messages, max_tokens=8192)
        if result and len(result) > 200:
            return result, 1  # 整份处理无法确定页数，标记为1
    except Exception:
        pass

    # 方法 3: pdfminer 提取文字
    try:
        from pdfminer.high_level import extract_text as pdfminer_extract
        text = pdfminer_extract(str(pdf_path))
        if text and len(text) > 100:
            return text, 1
    except ImportError:
        pass
    except Exception:
        pass

    # 方法 4: PyPDF2
    try:
        from PyPDF2 import PdfReader
        reader = PdfReader(str(pdf_path))
        pages = []
        for page in reader.pages:
            t = page.extract_text()
            if t:
                pages.append(t)
        text = "\n\n".join(pages)
        if text and len(text) > 100:
            return text, len(reader.pages)
    except ImportError:
        pass
    except Exception:
        pass

    return f"[无法提取 PDF 文本: {pdf_path.name}，请安装 pdf2image、pdfminer 或 PyPDF2]", 0


# ── 结构化抽取 ────────────────────────────────────────

EXTRACTION_PROMPT_TEMPLATE = """基于以下文档内容，提取关键事实信息。**只提取文档中明确写出的客观事实，不要做任何推理或判断**。严格按照 JSON 格式输出。

JSON schema（字段说明）：
- company_name: 公司名称（文档中出现的公司全称）
- industry: 所属行业（文档中明确提到的）
- sub_industry: 细分行业（文档中明确提到的，必须尽可能具体。如"耐辐照芯片"而非"半导体"或"模拟芯片"，"商业航天发射服务"而非"航天"或"军工"。如果文档提到多个细分领域，用顿号分隔列出所有细分领域）
- founding_year: 成立年份（文档中明确提到的数字，无法确定填 null）
- headquarters: 总部所在地（文档中明确提到的）
- product_maturity: 产品成熟度，从以下选项中选择一个最符合的：概念/原型/小批量/量产
  **严格判断标准**：
  - "量产"：文档明确提到已有批量交付、稳定出货、销售/订单/客户落地。BP里的"量产样机""工程样机""发布样机""可量产"都不能直接判成量产。
  - "小批量"：文档明确提到样机、工程样机、试产、小范围试用、验证机，但没有明确批量商业交付。
  - "原型"：有技术展示但无真实产品交付，或只有实验室/演示原型。
  - "概念"：只有PPT/构想，连原型都没有。
- valuation_hint: 估值信息（文档中明确提到的数字、投后估值、融资金额、释放股份，无则填"未披露"）
- business_model: 商业模式简述（100字内，只复述文档内容）
- product_service: 数组，核心产品或服务名称（优先提取明确的产品名、平台名、解决方案名，不要漏掉系列产品）
- target_market: 目标市场/应用场景（文档中明确提到的行业、客户群、使用场景）
- team_highlights: 数组，实际创始人/CEO/CTO/COO/管理层的姓名和职务（格式"姓名 - 职务/头衔 - 关键背景"）。只放公司内部全职管理层，不放外部顾问。
- advisors: 数组，科学顾问/专家委员会/外部顾问的姓名和职务（格式"姓名 - 职务/头衔 - 关键背景"）。**严格限制**：只有文档中明确标注了"科学顾问""顾问委员会""专家委员会""外部顾问"等字样，并且列出了具体人名时，才填入此字段。不要从"专家团队""行业专家"等模糊表述中推断顾问。如果文档没有明确的顾问名单，必须填空数组 []。
- financial_highlights: 对象，含 revenue/growth_rate/profitability/key_metrics（只填文档中明确出现的数据，无则填"未披露"）
- revenue_breakdown: 对象，营收按产品线/业务线分解（如文档提到"抗辐照芯片占65%收入"等，按品类拆分。格式如 {"抗辐照芯片": "65%", "车规级电源管理芯片": "20%", "检测服务": "15%"}。如果文档未提供分解，填 {}）
- competitive_advantages: 数组，文档中声称的竞争优势、技术亮点、壁垒
- risks: 数组，文档中提到的风险、挑战、限制、监管要求；如文档未提到则填 ["未提及"]
- use_of_proceeds: 资金用途/融资用途（文档中明确提到的）
- investor_highlights: 数组，已有投资方名称（文档中提到的，无则空数组）
- raw_financing_stage_text: 文档中关于融资阶段/轮次的原文表述（照抄，如"种子轮""天使轮""Pre-A"等。如果文档中没有提到任何融资阶段信息，填"未提及"）
- summary_100words: 100字以内的项目概述

**注意**：
1. financing_stage 字段不要在这里判断，由系统根据提取的事实自动推断
2. 只提取文档中明确出现的信息，不要推测或补充
3. 数字类信息必须精确，不要四舍五入或近似
4. team_highlights 只放实际管理层（创始人/CEO/CTO/COO/VP 等），advisors 放外部顾问。如果文档标题写了"科学顾问""顾问委员会"，对应人物必须放 advisors 而不是 team_highlights。每个人逐个保留姓名、职务、关键背景。
5. 不要因为文档出现"行业前景/应用前景/市场规模"就漏掉 target_market、competitive_advantages、product_service；这些字段宁可多列，也不要空着。
6. 如果文档出现"量产样机/工程样机/发布样机"，product_maturity 优先判为"小批量"而不是"量产"。

直接输出 JSON 对象，用 ```json``` 包裹。

文档内容：
{content}
"""


def _infer_financing_stage(profile: dict) -> tuple[str, str]:
    """基于提取的客观事实，用严格规则推断融资阶段。

    返回 (financing_stage, rationale)

    规则优先级：
    1. 文档明确提到融资阶段 → 优先采纳
    2. 否则基于产品成熟度+营收+团队做保守推断
    """
    raw_text = profile.get("raw_financing_stage_text", "未提及")
    product_maturity = profile.get("product_maturity", "")
    financials = profile.get("financial_highlights", {})
    revenue = financials.get("revenue", "未披露")
    team = profile.get("team_highlights", [])

    has_revenue = revenue not in ("未披露", "", "0", "无", None, "零营收")
    is_prototype = product_maturity in ("概念", "原型", "")
    is_small_batch = product_maturity == "小批量"
    is_mass = product_maturity == "量产"
    has_team = len(team) >= 2

    # BP 常把"量产样机/工程样机"写成量产，这里保守降级
    if is_mass and not has_revenue:
        is_mass = False
        is_small_batch = True

    if raw_text and raw_text != "未提及":
        stage = raw_text
        rationale = f"文档明确写了融资阶段：{raw_text}。"
        return stage, rationale

    if is_prototype and not has_revenue:
        if has_team:
            return "种子轮", "产品仍在概念/原型阶段，暂无营收，但已组建核心团队，符合种子轮特征。"
        return "种子轮", "产品仍在概念/原型阶段，暂无营收，符合种子轮特征。"

    if is_small_batch and not has_revenue:
        return "种子轮", "产品处于样机/小批量验证阶段，但无稳定收入或商业化交付证据，按种子轮保守判断。"

    if is_small_batch and has_revenue:
        return "天使轮", "产品已有样机或小批量验证，且出现初步收入，符合天使轮特征。"

    if is_mass and not has_revenue:
        return "天使轮", "产品已接近量产，但缺少稳定收入数据，按天使轮保守判断。"

    if is_mass and has_revenue:
        return "Pre-A轮", "产品已有量产和收入，处于从验证走向扩张的早期商业化阶段。"

    if has_revenue:
        return "天使轮", "出现收入信息，但产品成熟度信息不足，按天使轮保守判断。"

    return "种子轮", "缺少明确融资阶段和商业化证据，按种子轮保守判断。"


def extract_structured_info(ocr_text: str) -> dict[str, Any]:
    """用 VL 模型从 OCR 文本中抽取客观事实信息，融资阶段由规则推断"""
    # 截断过长文本（VL 模型上下文窗口足够大，尽量多喂）
    content = ocr_text[:32000]
    # ⚠️ 不能用 str.format()——OCR 文本中的花括号会被误解析为占位符导致 KeyError
    # 改用字符串拼接，安全地插入 content
    prompt = EXTRACTION_PROMPT_TEMPLATE.replace("{content}", content, 1)
    
    messages = [{"role": "user", "content": prompt}]
    response = _vl_chat(messages, max_tokens=4096)
    
    # 尝试解析 JSON
    profile = None
    try:
        # 找到 JSON 块
        if "```json" in response:
            json_str = response.split("```json")[1].split("```")[0].strip()
        elif "```" in response:
            json_str = response.split("```")[1].split("```")[0].strip()
        elif "{" in response:
            start = response.index("{")
            end = response.rindex("}") + 1
            json_str = response[start:end]
        else:
            json_str = response
        
        profile = json.loads(json_str)
    except (json.JSONDecodeError, ValueError):
        # JSON 解析失败，返回原始响应作为 summary
        return {
            "company_name": "",
            "industry": "",
            "extraction_error": f"JSON 解析失败: {str(response[:200])}",
            "summary_100words": response[:200],
        }
    
    # 用规则推断融资阶段（不依赖 VL 判断）
    if profile is not None:
        stage, rationale = _infer_financing_stage(profile)
        profile["financing_stage"] = stage
        profile["financing_stage_rationale"] = rationale

        # 后处理：验证提取结果不是幻觉
        profile = _validate_extracted_profile(profile, ocr_text)

    return profile or {"extraction_error": "空结果"}


OCR_FAILURE_MARKERS = (
    "无法提取 PDF 文本",
    "OCR失败",
    "VL识别无文字",
    "无文字内容",
    "请安装 pdf2image",
)
MIN_OCR_USEFUL_CHARS = 500
MIN_OCR_SUCCESS_RATIO = 0.6


def _ocr_healthcheck() -> dict[str, Any]:
    vl_api_base, vl_api_key, vl_model = _resolve_vl_config()
    python_packages = {
        "pdf2image": importlib.util.find_spec("pdf2image") is not None,
        "pdfminer": importlib.util.find_spec("pdfminer") is not None,
        "PyPDF2": importlib.util.find_spec("PyPDF2") is not None,
        "pptx": importlib.util.find_spec("pptx") is not None,
        "docx": importlib.util.find_spec("docx") is not None,
        "PIL": importlib.util.find_spec("PIL") is not None,
        "pytesseract": importlib.util.find_spec("pytesseract") is not None,
    }
    binaries = {
        "soffice": bool(shutil.which("soffice") or shutil.which("libreoffice") or Path("/Applications/LibreOffice.app/Contents/MacOS/soffice").exists()),
        "pdftoppm": bool(shutil.which("pdftoppm")),
    }
    return {
        "vl": {
            "base_url_present": bool(vl_api_base),
            "api_key_present": bool(vl_api_key),
            "model": vl_model,
        },
        "python_packages": python_packages,
        "binaries": binaries,
    }


def _line_is_ocr_noise(line: str) -> bool:
    stripped = line.strip()
    if not stripped:
        return True
    if re.fullmatch(r"[-—\s第0-9页/SlidePage?]+", stripped, re.I):
        return True
    if any(marker in stripped for marker in OCR_FAILURE_MARKERS):
        return True
    if stripped.startswith("[") and stripped.endswith("]"):
        return True
    return False


def _useful_text_from_ocr(ocr_text: str) -> str:
    useful_lines = [line.strip() for line in (ocr_text or "").splitlines() if not _line_is_ocr_noise(line)]
    return "\n".join(useful_lines).strip(" []（）():：，。,./\n\t-")


def _page_status(page_text: str) -> str:
    useful = _useful_text_from_ocr(page_text)
    if len(useful) >= 80:
        return "success"
    if any(marker in page_text for marker in OCR_FAILURE_MARKERS) or not useful:
        return "failed"
    return "empty"


def _ocr_page_records(ocr_text: str, pages_processed: int) -> list[dict[str, Any]]:
    text = ocr_text or ""
    matches = list(re.finditer(r"^---\s*(?:第)?([^\n-]+?)(?:页|Page|Slide)?\s*---\s*$", text, re.M | re.I))
    records: list[dict[str, Any]] = []
    if matches:
        for idx, match in enumerate(matches):
            start = match.end()
            end = matches[idx + 1].start() if idx + 1 < len(matches) else len(text)
            body = text[start:end].strip()
            useful = _useful_text_from_ocr(body)
            records.append({
                "page_index": idx + 1,
                "label": match.group(1).strip(),
                "status": _page_status(body),
                "chars": len(body),
                "useful_chars": len(useful),
            })
    elif text:
        useful = _useful_text_from_ocr(text)
        records.append({
            "page_index": 1,
            "label": "document",
            "status": _page_status(text),
            "chars": len(text),
            "useful_chars": len(useful),
        })
    if matches:
        while len(records) < pages_processed:
            records.append({
                "page_index": len(records) + 1,
                "label": str(len(records) + 1),
                "status": "failed",
                "chars": 0,
                "useful_chars": 0,
            })
    return records


def _build_ocr_manifest(input_path: Path, ext: str, ocr_text: str, pages_processed: int) -> dict[str, Any]:
    pages = _ocr_page_records(ocr_text, pages_processed)
    success_pages = sum(1 for page in pages if page.get("status") == "success")
    has_explicit_page_markers = bool(re.search(r"^---\s*(?:第)?[^\n-]+?(?:页|Page|Slide)?\s*---\s*$", ocr_text or "", re.M | re.I))
    plain_full_document = bool((ocr_text or "").strip()) and not has_explicit_page_markers and len(pages) == 1 and pages[0].get("status") == "success"
    total_pages = 1 if plain_full_document else max(len(pages), pages_processed, 0)
    useful_chars = len(_useful_text_from_ocr(ocr_text))
    success_ratio = 1.0 if plain_full_document else (round(success_pages / total_pages, 4) if total_pages else 0.0)
    reason = ""
    verdict = "PASS"
    if not (ocr_text or "").strip() or len(ocr_text or "") < 50:
        verdict = "FAIL"
        reason = "OCR_EMPTY_OR_TOO_SHORT"
    elif useful_chars < MIN_OCR_USEFUL_CHARS and any(marker in (ocr_text or "") for marker in OCR_FAILURE_MARKERS):
        verdict = "FAIL"
        reason = "OCR_RESULT_UNUSABLE"
    elif pages_processed <= 0:
        verdict = "FAIL"
        reason = "OCR_NO_PAGES_PROCESSED"
    elif total_pages >= 3 and success_ratio < MIN_OCR_SUCCESS_RATIO:
        verdict = "FAIL"
        reason = "OCR_LOW_SUCCESS_RATIO"
    return {
        "schema_version": "bp_ocr_manifest.v1",
        "input_file": str(input_path),
        "input_type": ext,
        "generated_at": time.strftime("%Y-%m-%dT%H:%M:%S"),
        "healthcheck": _ocr_healthcheck(),
        "quality_gate": {
            "verdict": verdict,
            "reason": reason,
            "pages_processed": pages_processed,
            "page_count": total_pages,
            "success_pages": success_pages,
            "success_ratio": success_ratio,
            "ocr_text_length": len(ocr_text or ""),
            "useful_chars": useful_chars,
            "min_success_ratio": MIN_OCR_SUCCESS_RATIO,
            "min_useful_chars": MIN_OCR_USEFUL_CHARS,
        },
        "pages": pages,
    }


def _ocr_text_unusable_reason(ocr_text: str, pages_processed: int) -> str:
    manifest = _build_ocr_manifest(Path(""), "", ocr_text, pages_processed)
    return manifest.get("quality_gate", {}).get("reason", "") if manifest.get("quality_gate", {}).get("verdict") == "FAIL" else ""


def _validate_extracted_profile(profile: dict, ocr_text: str) -> dict:
    """验证 VL 提取结果不是幻觉。对每个提取字段做交叉检查。"""
    import re as _re

    # 1. 验证 advisors：提取的人名必须在 OCR 文本中出现
    advisors = profile.get("advisors") or []
    if advisors:
        validated_advisors = []
        for advisor in advisors:
            name = str(advisor).split(" - ")[0].split("-")[0].strip()
            if name and len(name) >= 2 and name in ocr_text:
                validated_advisors.append(advisor)
            elif name:
                print(f"  ⚠️ 顾问 '{name}' 未在 OCR 文本中找到，可能是幻觉，已过滤", flush=True)
        if len(validated_advisors) != len(advisors):
            print(f"  📋 顾问验证: {len(advisors)} → {len(validated_advisors)}（过滤了幻觉项）", flush=True)
        profile["advisors"] = validated_advisors

    # 2. 验证 team_highlights：提取的人名必须在 OCR 文本中出现
    team = profile.get("team_highlights") or []
    if team:
        validated_team = []
        for member in team:
            name = str(member).split(" - ")[0].split("-")[0].strip()
            if name and len(name) >= 2 and name in ocr_text:
                validated_team.append(member)
            elif name:
                print(f"  ⚠️ 团队成员 '{name}' 未在 OCR 文本中找到，可能是幻觉，已过滤", flush=True)
        if len(validated_team) != len(team):
            print(f"  📋 团队验证: {len(team)} → {len(validated_team)}（过滤了幻觉项）", flush=True)
        profile["team_highlights"] = validated_team

    # 3. 验证 company_name：必须在 OCR 文本中出现
    company_name = profile.get("company_name", "")
    if company_name and company_name not in ocr_text:
        # 尝试模糊匹配（去掉"有限公司"等后缀）
        short_name = company_name.replace("有限公司", "").replace("股份", "").replace("科技", "").strip()
        if short_name and short_name not in ocr_text:
            print(f"  ⚠️ 公司名 '{company_name}' 未在 OCR 文本中找到，可能是幻觉", flush=True)

    return profile


# ── 主入口 ────────────────────────────────────────────

def run_document_intake(job_ctx, input_file: str) -> dict[str, Any]:
    """BP 文档入库主流程

    Args:
        job_ctx: JobContext（含 workspace）
        input_file: 输入文件路径

    Returns:
        标准结果字典
    """
    from runtime.profiles.base import JobContext

    input_path = Path(input_file)
    if not input_path.exists():
        return {
            "ok": False,
            "mode": "legacy_wrapped",
            "phase": "phase01_document_intake",
            "job_id": job_ctx.job_id,
            "error": f"输入文件不存在: {input_file}",
        }

    # 确定输出目录
    workspace = getattr(job_ctx, "workspace", None)
    if workspace is not None:
        output_dir = workspace.root
        extraction_dir = workspace.extraction_dir
    else:
        # 兼容旧路径
        from pathlib import Path as P
        rt = P(__file__).resolve().parents[2]
        output_dir = rt / "tasks" / job_ctx.job_id
        extraction_dir = output_dir
        output_dir.mkdir(parents=True, exist_ok=True)

    ext = input_path.suffix.lower()
    if ext not in SUPPORTED_EXTENSIONS:
        return {
            "ok": False,
            "mode": "legacy_wrapped",
            "phase": "phase01_document_intake",
            "job_id": job_ctx.job_id,
            "error": f"不支持的文件类型: {ext}，支持: {SUPPORTED_EXTENSIONS}",
        }

    # 1. OCR 提取文字
    ocr_text = ""
    pages_processed = 0

    try:
        if ext in (".png", ".jpg", ".jpeg", ".bmp", ".tiff", ".webp"):
            # 单张图片
            ocr_text = _ocr_image(input_path)
            pages_processed = 1

        elif ext == ".pdf":
            ocr_text, pages_processed = _ocr_pdf(input_path, extraction_dir)

        elif ext in (".pptx", ".ppt"):
            # PPTX: 优先 LibreOffice 渲染整页 → VL OCR（覆盖图片型页面）
            #       fallback: python-pptx 读文字 + 嵌入图片 VL OCR
            page_texts = []
            libreoffice_ok = False

            # ── Step 0: LibreOffice 渲染整页为 PNG ──
            try:
                print("  📄 正在用 LibreOffice 转换 PPTX → PDF → PNG ...", flush=True)
                rendered_dir = _render_pptx_with_libreoffice(input_path, extraction_dir)
                if rendered_dir and any(rendered_dir.glob("*.png")):
                    libreoffice_ok = True
                    slide_images = sorted(rendered_dir.glob("*.png"))
                    total_slides = len(slide_images)
                    print(f"  ✅ 渲染完成，共 {total_slides} 页，开始逐页 OCR ...", flush=True)
                    for img_path in slide_images:
                        slide_idx = img_path.stem.split("-")[-1] if "-" in img_path.stem else ""
                        try:
                            t = _ocr_image(img_path, page_hint=f"第{slide_idx}页" if slide_idx else "")
                            if t and len(t) > 10:
                                page_texts.append(f"--- 第{slide_idx or '?'}页 ---\n{t}")
                            else:
                                page_texts.append(f"--- 第{slide_idx or '?'}页 ---\n（VL识别无文字）")
                            pages_processed += 1
                            if pages_processed % 5 == 0 or pages_processed == total_slides:
                                print(f"    OCR 进度: {pages_processed}/{total_slides}", flush=True)
                        except Exception as e:
                            page_texts.append(f"--- 第{slide_idx or '?'}页 ---\n[OCR失败: {e}]")
                    print(f"  📊 LibreOffice 渲染 + VL OCR: {len(slide_images)} 页", flush=True)
                else:
                    print("  ⚠ LibreOffice 渲染未产出 PNG，回退到 python-pptx", flush=True)
            except Exception as e:
                print(f"  ⚠ LibreOffice 渲染失败 ({e})，回退到 python-pptx")

            # ── Fallback: python-pptx 读文字 + 嵌入图片 VL OCR ──
            if not libreoffice_ok:
                try:
                    from pptx import Presentation
                    prs = Presentation(str(input_path))
                    img_dir = extraction_dir / "slide_images"
                    img_dir.mkdir(parents=True, exist_ok=True)

                    for i, slide in enumerate(prs.slides, 1):
                        texts = []
                        for shape in slide.shapes:
                            if shape.has_text_frame:
                                for para in shape.text_frame.paragraphs:
                                    t = para.text.strip()
                                    if t:
                                        texts.append(t)
                        if texts:
                            page_texts.append(f"--- 第{i}页 ---\n" + "\n".join(texts))
                        else:
                            page_texts.append(f"--- 第{i}页 ---\n（无文字内容）")
                        pages_processed += 1

                    # 补充：嵌入图片 VL OCR
                    images = _extract_images_from_pptx(input_path, img_dir)
                    if images:
                        for idx, img in enumerate(images[:30], 1):
                            try:
                                img_ocr = _ocr_image(img, page_hint=f"嵌入图片{idx}")
                                if img_ocr and len(img_ocr) > 20:
                                    page_texts.append(f"\n--- 嵌入图片{idx} OCR ---\n{img_ocr}")
                            except Exception as e:
                                page_texts.append(f"\n--- 嵌入图片{idx} ---\n[OCR失败: {e}]")

                except ImportError:
                    img_dir = extraction_dir / "slide_images"
                    images = _extract_images_from_pptx(input_path, img_dir)
                    for i, img in enumerate(images, 1):
                        try:
                            t = _ocr_image(img, page_hint=f"第{i}页/幻灯片")
                            page_texts.append(f"--- 第{i}页 ---\n{t}")
                            pages_processed += 1
                        except Exception as e:
                            page_texts.append(f"--- 第{i}页 ---\n[OCR失败: {e}]")
            ocr_text = "\n\n".join(page_texts)

        elif ext in (".docx", ".doc"):
            # DOCX: 优先用 python-docx 读文字，再对图片做 VL OCR
            page_texts = []
            
            # Step 1: python-docx 读文字
            doc_text = ""
            try:
                from docx import Document
                doc = Document(str(input_path))
                doc_text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
                # 也读表格
                for table in doc.tables:
                    for row in table.rows:
                        row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                        if row_text:
                            doc_text += "\n" + row_text
            except ImportError:
                pass
            except Exception:
                pass

            if doc_text:
                page_texts.append(f"--- 文档文字 ---\n{doc_text}")
                pages_processed += 1

            # Step 2: OCR 嵌入图片（补充图表信息）
            img_dir = extraction_dir / "doc_images"
            images = _extract_images_from_docx(input_path, img_dir)
            for i, img in enumerate(images, 1):
                try:
                    t = _ocr_image(img, page_hint=f"图片{i}")
                    page_texts.append(f"--- 图片{i} ---\n{t}")
                    pages_processed += 1
                except Exception as e:
                    page_texts.append(f"--- 图片{i} ---\n[OCR失败: {e}]")
            ocr_text = "\n\n".join(page_texts)

    except Exception as e:
        return {
            "ok": False,
            "mode": "legacy_wrapped",
            "phase": "phase01_document_intake",
            "job_id": job_ctx.job_id,
            "error": f"OCR 处理失败: {e}",
        }

    manifest = _build_ocr_manifest(input_path, ext, ocr_text, pages_processed)
    manifest_path = output_dir / "bp_ocr_manifest.json"
    manifest_path.write_text(json.dumps(manifest, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    unusable_reason = manifest.get("quality_gate", {}).get("reason", "") if manifest.get("quality_gate", {}).get("verdict") == "FAIL" else ""
    if unusable_reason:
        return {
            "ok": False,
            "mode": "legacy_wrapped",
            "phase": "phase01_document_intake",
            "job_id": job_ctx.job_id,
            "error": unusable_reason,
            "details": {
                "pages_processed": pages_processed,
                "ocr_text_length": len(ocr_text or ""),
                "manifest_path": str(manifest_path),
            },
        }

    # 2. 写 OCR 文本
    ocr_path = output_dir / "bp_ocr_text.txt"
    ocr_path.write_text(ocr_text, encoding="utf-8")

    # 3. 结构化抽取
    profile = {}
    try:
        profile = extract_structured_info(ocr_text)
    except Exception as e:
        profile = {"extraction_error": str(e)}

    profile_path = output_dir / "bp_step0_profile.json"
    profile_path.write_text(
        json.dumps(profile, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )

    # 同时写一个 Markdown 版本方便子代理阅读
    profile_md_path = output_dir / "bp_step0_profile.md"
    md_lines = ["# BP Step0 Profile\n"]
    for k, v in profile.items():
        if isinstance(v, list):
            md_lines.append(f"## {k}")
            for item in v:
                md_lines.append(f"- {item}")
        elif isinstance(v, dict):
            md_lines.append(f"## {k}")
            for sk, sv in v.items():
                md_lines.append(f"- **{sk}**: {sv}")
        else:
            md_lines.append(f"- **{k}**: {v}")
    profile_md_path.write_text("\n".join(md_lines), encoding="utf-8")

    # 4. 记录产物
    try:
        from runtime.orchestrator.state_store import StateStore
        ss = StateStore(Path(__file__).resolve().parents[2])
        ss.record_artifact(job_ctx.job_id, "bp_ocr_text", ocr_path)
        ss.record_artifact(job_ctx.job_id, "bp_ocr_manifest", manifest_path)
        ss.record_artifact(job_ctx.job_id, "bp_step0_profile", profile_path)
    except Exception:
        pass

    return {
        "ok": True,
        "mode": "legacy_wrapped",
        "phase": "phase01_document_intake",
        "job_id": job_ctx.job_id,
        "result": {
            "input_file": str(input_path),
            "input_type": ext,
            "pages_processed": pages_processed,
            "ocr_text_length": len(ocr_text),
            "ocr_path": str(ocr_path),
            "profile_path": str(profile_path),
            "company_name": profile.get("company_name", ""),
            "industry": profile.get("industry", ""),
            "financing_stage": profile.get("financing_stage", ""),
        },
    }
