#!/usr/bin/env python3
"""
BP OCR — 支持 PDF 和 PPTX 两种格式的 BP 文档 OCR 提取
"""
from __future__ import annotations
import argparse
import os
import sys
import time
import io
import base64
from pathlib import Path

try:
    from openai import OpenAI
except ImportError:  # pragma: no cover - exercised only when dependency is missing at runtime
    OpenAI = None

ROOT = Path(__file__).resolve().parent.parent
DEFAULT_BP_OCR_MODEL = 'qwen3-vl-30b-a3b-instruct'
DEFAULT_BP_OCR_BASE_URL = ''  # Set via BP_OCR_BASE_URL env var

# ═══════════════════════════════════════════════
# PPTX 转图片
# ═══════════════════════════════════════════════
def _pptx_to_images_pillow(pptx_path: str, dpi: int = 300) -> list:
    """
    将 PPTX 每一页转为图片。
    用 python-pptx 提取 shape 文本 + 可选的图片渲染。
    """
    try:
        from pptx import Presentation
        from pptx.util import Emu
        from PIL import Image, ImageDraw, ImageFont
    except ImportError:
        raise ImportError("需要安装 python-pptx: pip3 install python-pptx")

    prs = Presentation(pptx_path)
    images = []

    for slide_idx, slide in enumerate(prs.slides):
        # 计算幻灯片尺寸（单位：EMU → 像素）
        slide_width = int(prs.slide_width / 914400 * dpi / 96)
        slide_height = int(prs.slide_height / 914400 * dpi / 96)
        if slide_width < 100: slide_width = 1920
        if slide_height < 100: slide_height = 1080

        # 创建白色背景
        img = Image.new('RGB', (slide_width, slide_height), 'white')
        draw = ImageDraw.Draw(img)

        try:
            font = ImageFont.truetype("/System/Library/Fonts/PingFang.ttc", max(12, slide_width // 80))
        except:
            try:
                font = ImageFont.truetype("/System/Library/Fonts/STHeiti Medium.ttc", max(12, slide_width // 80))
            except:
                font = ImageFont.load_default()

        y_pos = 20
        line_height = font.size + 6

        # 提取所有文本
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        # 限制行宽
                        max_chars = slide_width // (font.size // 2)
                        for line_idx in range(0, max(1, (len(text) + max_chars - 1) // max_chars)):
                            if y_pos + line_height > slide_height - 20:
                                break
                            chunk = text[line_idx * max_chars:(line_idx + 1) * max_chars]
                            draw.text((20, y_pos), chunk, fill=(50, 50, 50), font=font)
                            y_pos += line_height

            # 如果有图片 shape
            if shape.size_rel is not None or hasattr(shape, 'image'):
                try:
                    if hasattr(shape, 'image') and shape.image is not None:
                        img_obj = shape.image
                        slide_img = Image.open(io.BytesIO(img_obj.blob))
                        slide_img.thumbnail((slide_width, slide_height))
                        img.paste(slide_img, (0, 0))
                except:
                    pass

        img_byte = io.BytesIO()
        img.save(img_byte, format='JPEG', quality=90)
        images.append(img_byte)

    return images


def _pptx_to_text_direct(pptx_path: str) -> str:
    """直接从 PPTX 提取纯文本（快，但可能漏掉纯图片幻灯片）。"""
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError("需要安装 python-pptx: pip3 install python-pptx")

    prs = Presentation(pptx_path)
    lines = []
    for slide_idx, slide in enumerate(prs.slides):
        lines.append(f"\n--- Slide {slide_idx + 1} ---\n")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        lines.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    lines.append(" | ".join(cells))
    return '\n'.join(lines)


# ═══════════════════════════════════════════════
# PDF OCR（原始逻辑）
# ═══════════════════════════════════════════════
def _pdf_to_images(pdf_path: str, dpi: int = 200) -> list:
    """用 pdf2image 把 PDF 每页转成图片。"""
    from pdf2image import convert_from_path
    images = convert_from_path(pdf_path, dpi=dpi, thread_count=4)
    return images


def _ocr_images_with_qwen_vl(
    images_or_bytes: list,
    api_key: str,
    model: str = DEFAULT_BP_OCR_MODEL,
    base_url: str = DEFAULT_BP_OCR_BASE_URL,
) -> str:
    """用 OpenAI-compatible 视觉模型 OCR 识别图片中的文字。"""
    if OpenAI is None:
        raise ImportError("需要安装 openai: pip install openai")

    client = OpenAI(api_key=api_key, base_url=base_url)
    full_text = []
    total = len(images_or_bytes)

    for i, img_source in enumerate(images_or_bytes):
        if hasattr(img_source, 'getvalue'):
            # BytesIO 对象 → 转 base64
            img_byte = img_source.getvalue()
            img_base64 = base64.b64encode(img_byte).decode()
        elif isinstance(img_source, bytes):
            img_base64 = base64.b64encode(img_source).decode()
        elif hasattr(img_source, 'save'):
            # PIL Image 对象
            img_byte = io.BytesIO()
            img_source.save(img_byte, format='JPEG', quality=85)
            img_base64 = base64.b64encode(img_byte.getvalue()).decode()
        else:
            # 文件路径
            with open(img_source, 'rb') as f:
                img_base64 = base64.b64encode(f.read()).decode()

        messages = [{
            'role': 'user',
            'content': [
                {
                    'type': 'image_url',
                    'image_url': {'url': f'data:image/jpeg;base64,{img_base64}'},
                },
                {
                    'type': 'text',
                    'text': '请识别图片中的所有文字，保持原有格式，直接输出文字内容不要解释。',
                },
            ],
        }]

        try:
            response = client.chat.completions.create(
                model=model,
                messages=messages,
                temperature=0.1,
            )
            text = response.choices[0].message.content if response.choices else ''
            if isinstance(text, list):
                text = ''.join(item.get('text', '') if isinstance(item, dict) else str(item) for item in text)
            text = text or ''
            full_text.append(f"\n--- Page {i+1} ---\n{text}")
            print(f"  ✅ 第 {i+1}/{total} 页 ({len(text)} 字)")

        except Exception as e:
            print(f"  ⚠ 第 {i+1} 页异常: {e}", file=sys.stderr)
            break

        time.sleep(0.3)

    return '\n'.join(full_text)


def _ocr_tesseract_fallback(pdf_path: str) -> str:
    """Tesseract OCR 降级。"""
    from pdf2image import convert_from_path
    try:
        import pytesseract
    except ImportError:
        raise ImportError("需要安装 pytesseract: pip3 install pytesseract")

    images = convert_from_path(pdf_path, dpi=300, thread_count=4)
    full_text = []
    for i, img in enumerate(images):
        text = pytesseract.image_to_string(img, lang='chi_sim+eng')
        full_text.append(text)
    return '\n'.join(full_text)


# ═══════════════════════════════════════════════
# 主入口
# ═══════════════════════════════════════════════
def _read_env_file() -> dict[str, str]:
    """读取项目 .credentials 中的 OCR/LLM 配置。"""
    cred_path = ROOT / '.credentials' / 'investment-research.env'
    values: dict[str, str] = {}
    if not cred_path.exists():
        return values
    for raw_line in cred_path.read_text(encoding='utf-8').splitlines():
        line = raw_line.strip()
        if not line or line.startswith('#') or '=' not in line:
            continue
        key, value = line.split('=', 1)
        values[key.strip()] = value.strip().strip("'\"")
    return values


def load_ocr_config() -> dict[str, str]:
    """加载 BP OCR 的 OpenAI-compatible API 配置。"""
    file_values = _read_env_file()
    api_key = os.environ.get('BP_OCR_API_KEY') or os.environ.get('DASHSCOPE_API_KEY') or file_values.get('BP_OCR_API_KEY') or file_values.get('DASHSCOPE_API_KEY')
    if not api_key:
        raise EnvironmentError("找不到 BP_OCR_API_KEY 或 DASHSCOPE_API_KEY")
    return {
        'api_key': api_key,
        'model': os.environ.get('BP_OCR_MODEL') or file_values.get('BP_OCR_MODEL') or DEFAULT_BP_OCR_MODEL,
        'base_url': os.environ.get('BP_OCR_BASE_URL') or file_values.get('BP_OCR_BASE_URL') or DEFAULT_BP_OCR_BASE_URL,
    }


def load_api_key() -> str:
    """兼容旧调用：返回 BP OCR API Key。"""
    return load_ocr_config()['api_key']


def mode_auto(file_path: str) -> str:
    """
    自动识别文件类型（PDF 或 PPTX）并提取文本。
    优先用 qwen-vl OCR，fallback Tesseract。
    """
    ext_path = Path(file_path).suffix.lower()

    # === PPTX 路径 ===
    if ext_path == '.pptx':
        print(f"  📄 检测到 PPTX 文件: {file_path}")

        # 方法 1: 直接提取文本（最快，但可能漏图）
        print("  📝 提取 PPTX 文本内容...")
        direct_text = _pptx_to_text_direct(file_path)

        # 方法 2: 幻灯片 OCR（获取图片/图表内容）
        print(f"  🔍 PPTX 幻灯片 OCR（共 {len(_pptx_to_images_pillow(file_path))} 页）...")
        images = _pptx_to_images_pillow(file_path, dpi=200)
        ocr_config = load_ocr_config()
        ocr_text = _ocr_images_with_qwen_vl(images, **ocr_config)

        # 合并：直接文本 + OCR 文本（去重）
        return f"[PPTX 文本提取]\n{direct_text}\n\n[PPTX 幻灯片 OCR]\n{ocr_text}"

    # === PDF 路径 ===
    elif ext_path == '.pdf':
        print(f"  📄 检测到 PDF 文件: {file_path}")
        try:
            ocr_config = load_ocr_config()
            images = _pdf_to_images(file_path, dpi=200)
            print(f"  🔍 qwen3-vl OCR: {len(images)} 页")
            return _ocr_images_with_qwen_vl(images, **ocr_config)
        except (EnvironmentError, ImportError) as e:
            print(f"  ⚠ qwen-vl 不可用，回退 Tesseract: {e}")
            return _ocr_tesseract_fallback(file_path)
        except Exception as e:
            print(f"  ⚠ OCR 异常，回退 Tesseract: {e}")
            return _ocr_tesseract_fallback(file_path)

    else:
        raise ValueError(f"不支持的文件格式: {ext_path}，仅支持 .pdf 和 .pptx")


if __name__ == '__main__':
    ap = argparse.ArgumentParser(description='BP OCR（支持 PDF 和 PPTX）')
    ap.add_argument('--pdf', '--file', default=None, help='BP 文档路径')
    ap.add_argument('--mode', default='auto', choices=['auto', 'qwen_vl', 'tesseract'])
    args = ap.parse_args()
    if not args.pdf:
        ap.print_help()
        sys.exit(1)
    text = mode_auto(args.pdf)
    print('\n' + text[:500] + ('...' if len(text) > 500 else ''))
